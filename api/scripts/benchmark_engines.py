"""Benchmark conversion engines on source-image fidelity.

This script compares multiple conversion paths on the same PDF pages:

- `local_tesseract` (main API pipeline with `parse_provider=local`, OCR provider=`tesseract`)
- `local_paddle_local` (main API pipeline with `parse_provider=local`, OCR provider=`paddle_local`)
- `local_paddle` (main API pipeline with `parse_provider=local`, OCR provider=`paddle`; requires API key)
- `local_aiocr` (main API pipeline with `parse_provider=local`, OCR provider=`aiocr`; requires API key)
- `v2` (legacy alias: `parse_provider=v2`; now routes to local+fullpage+AI OCR; requires API key)
- `v3_openai_chat` (isolated v3 pipeline, chat OCR backend)
- `v3_paddle_doc` (isolated v3 pipeline, Paddle doc_parser backend)

For each page, it generates a synthetic final image by drawing OCR text back
onto cleaned backgrounds, then compares against the original rendered PDF page.
Outputs include per-page artifacts and a JSON report with ranked engines.
The ranking combines image fidelity and editability proxy (PPT text content).

Example:
  env -u http_proxy -u https_proxy -u HTTP_PROXY -u HTTPS_PROXY -u ALL_PROXY \
    api/.venv/bin/python api/scripts/benchmark_engines.py \
      --pdf ./+AI智能体开发大学生指南.pdf \
      --pages 1-3 \
      --out-dir ./test/benchmark-engines \
      --api-key "$SILICONFLOW_API_KEY" \
      --base-url https://api.siliconflow.cn/v1 \
      --model PaddlePaddle/PaddleOCR-VL-1.5
"""

from __future__ import annotations

import argparse
import json
import os
import random
import shutil
import socket
import sys
import time
import traceback
import uuid
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

import cv2
import numpy as np
import pymupdf
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


def _ensure_api_on_path() -> Path:
    api_dir = Path(__file__).resolve().parents[1]
    if str(api_dir) not in sys.path:
        sys.path.insert(0, str(api_dir))
    return api_dir


def _ensure_v3_on_path(api_dir: Path) -> None:
    # Removed: v2/v3 packages deleted
    pass


def _clear_proxy_env() -> None:
    for key in (
        "http_proxy",
        "https_proxy",
        "HTTP_PROXY",
        "HTTPS_PROXY",
        "ALL_PROXY",
    ):
        os.environ.pop(key, None)


def _patch_dns(host: str | None, ip: str | None) -> Callable[[], None]:
    if not host or not ip:
        return lambda: None

    original = socket.getaddrinfo

    def patched(target_host: str, port: int, *args: Any, **kwargs: Any):
        if target_host == host:
            return original(ip, port, *args, **kwargs)
        return original(target_host, port, *args, **kwargs)

    socket.getaddrinfo = patched

    def restore() -> None:
        socket.getaddrinfo = original

    return restore


def _parse_pages_spec(
    spec: str, *, max_pages: int, seed: int | None = None
) -> list[int]:
    out: set[int] = set()
    tokens = [t.strip() for t in (spec or "").split(",") if t.strip()]
    for token in tokens:
        lowered = token.lower()
        if lowered == "all":
            for page_no in range(1, int(max_pages) + 1):
                out.add(int(page_no))
            continue

        if lowered.startswith(("random:", "rand:", "sample:")):
            _, raw_n = lowered.split(":", 1)
            try:
                n = int(raw_n)
            except Exception:
                n = 0
            if n <= 0:
                continue
            n = min(int(n), int(max_pages))
            if n <= 0:
                continue
            if seed is None:
                seed = int(time.time_ns() % (2**31 - 1))
            rng = random.Random(int(seed))
            choices = rng.sample(list(range(1, int(max_pages) + 1)), n)
            for page_no in choices:
                out.add(int(page_no))
            continue

        if "-" in token:
            a, b = token.split("-", 1)
            start = int(a)
            end = int(b)
            if start > end:
                start, end = end, start
            for page_no in range(start, end + 1):
                if 1 <= page_no <= max_pages:
                    out.add(page_no)
        else:
            page_no = int(token)
            if 1 <= page_no <= max_pages:
                out.add(page_no)
    return sorted(out)


def _group_contiguous_spans(page_numbers: list[int]) -> list[tuple[int, int]]:
    """Group page numbers into contiguous [start,end] spans (1-based)."""

    unique = sorted({int(p) for p in page_numbers if int(p) > 0})
    if not unique:
        return []

    spans: list[list[int]] = []
    for page_no in unique:
        if not spans or page_no != spans[-1][1] + 1:
            spans.append([page_no, page_no])
        else:
            spans[-1][1] = page_no

    return [(int(a), int(b)) for a, b in spans if a > 0 and b >= a]


def _ensure_dir(path: Path) -> Path:
    path.mkdir(parents=True, exist_ok=True)
    return path


def _render_source_page(
    pdf_path: Path,
    *,
    page_no: int,
    dpi: int,
    out_path: Path,
) -> Path:
    with pymupdf.open(str(pdf_path)) as doc:
        page = doc.load_page(page_no - 1)
        pix = page.get_pixmap(dpi=int(dpi), alpha=False)  # type: ignore[attr-defined]
        pix.save(str(out_path))
    return out_path


def _read_color_image(path: Path) -> np.ndarray:
    img = cv2.imread(str(path), cv2.IMREAD_COLOR)
    if img is None:
        raise RuntimeError(f"failed to read image: {path}")
    return img


def _align_image_like(ref: np.ndarray, img: np.ndarray) -> np.ndarray:
    h, w = ref.shape[:2]
    return cv2.resize(img, (int(w), int(h)), interpolation=cv2.INTER_AREA)


def _image_metrics(a: np.ndarray, b: np.ndarray) -> dict[str, float]:
    diff = cv2.absdiff(a, b)
    gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
    mae = float(np.mean(diff)) / 255.0
    rmse = (
        float(np.sqrt(np.mean((a.astype(np.float32) - b.astype(np.float32)) ** 2)))
        / 255.0
    )
    high_diff_ratio = float(np.mean((gray > 32).astype(np.float32)))
    return {
        "mae": mae,
        "rmse": rmse,
        "high_diff_ratio": high_diff_ratio,
    }


def _save_side_by_side(
    *,
    left: np.ndarray,
    right: np.ndarray,
    left_label: str,
    right_label: str,
    out_path: Path,
) -> None:
    h, w = left.shape[:2]
    bar = np.full((48, w * 2, 3), 255, dtype=np.uint8)
    cv2.putText(
        bar,
        left_label,
        (20, 32),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.8,
        (0, 0, 0),
        2,
        cv2.LINE_AA,
    )
    cv2.putText(
        bar,
        right_label,
        (w + 20, 32),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.8,
        (0, 0, 0),
        2,
        cv2.LINE_AA,
    )
    canvas = np.hstack([left, right])
    cv2.imwrite(str(out_path), np.vstack([bar, canvas]))


def _save_heat_overlay(*, base: np.ndarray, other: np.ndarray, out_path: Path) -> None:
    diff = cv2.absdiff(base, other)
    gray = cv2.cvtColor(diff, cv2.COLOR_BGR2GRAY)
    heat = cv2.applyColorMap(gray, cv2.COLORMAP_JET)
    mixed = cv2.addWeighted(base, 0.55, heat, 0.45, 0)
    cv2.imwrite(str(out_path), mixed)


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if not path.exists():
            continue
        try:
            return ImageFont.truetype(str(path), size=size)
        except Exception:
            continue
    return ImageFont.load_default()


def _normalize_line_items(raw_items: list[dict[str, Any]]) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for item in raw_items:
        if not isinstance(item, dict):
            continue
        text = str(item.get("text") or "").strip()
        bbox = item.get("bbox")
        if not text or not isinstance(bbox, list) or len(bbox) != 4:
            continue
        try:
            x0, y0, x1, y1 = [float(v) for v in bbox]
        except Exception:
            continue
        x0, x1 = min(x0, x1), max(x0, x1)
        y0, y1 = min(y0, y1), max(y0, y1)
        if x1 <= x0 or y1 <= y0:
            continue
        out.append({"text": text, "bbox": [x0, y0, x1, y1]})
    out.sort(key=lambda it: (it["bbox"][1], it["bbox"][0]))
    return out


def _synthesize_from_clean(
    *,
    clean_image: Path,
    lines: list[dict[str, Any]],
    out_path: Path,
) -> Path:
    image = Image.open(clean_image).convert("RGB")
    draw = ImageDraw.Draw(image)
    norm_lines = _normalize_line_items(lines)
    for line in norm_lines:
        text = line["text"]
        x0, y0, x1, y1 = [float(v) for v in line["bbox"]]
        h = max(1.0, y1 - y0)
        size = max(8, min(64, int(h * 0.78)))
        font = _load_font(size=size)
        draw.text((x0, y0), text, fill=(0, 0, 0), font=font)
    image.save(str(out_path))
    return out_path


def _find_first_existing(candidates: list[Path]) -> Path | None:
    for candidate in candidates:
        if candidate.exists():
            return candidate
    return None


def _load_local_lines_px(
    *,
    job_dir: Path,
    page_no: int,
    clean_image: Path,
) -> list[dict[str, Any]]:
    ir_path = _find_first_existing(
        [
            job_dir / "ir.ocr.json",
            job_dir / "ir.json",
            job_dir / "ir.parsed.json",
        ]
    )
    if ir_path is None:
        return []

    try:
        ir = json.loads(ir_path.read_text(encoding="utf-8"))
    except Exception:
        return []

    pages = ir.get("pages") or []
    if not isinstance(pages, list):
        return []

    page_obj: dict[str, Any] | None = None
    target_index = int(page_no - 1)
    for item in pages:
        if not isinstance(item, dict):
            continue
        idx = item.get("page_index")
        try:
            idx_int = int(idx) if idx is not None else None
        except Exception:
            idx_int = None
        if idx_int == target_index or idx_int == page_no:
            page_obj = item
            break

    if page_obj is None:
        return []

    try:
        with Image.open(clean_image) as im:
            width_px, height_px = im.size
    except Exception:
        return []

    page_w_pt = float(page_obj.get("page_width_pt") or 0.0)
    page_h_pt = float(page_obj.get("page_height_pt") or 0.0)
    if page_w_pt <= 0 or page_h_pt <= 0:
        return []

    lines: list[dict[str, Any]] = []
    for el in page_obj.get("elements") or []:
        if not isinstance(el, dict):
            continue
        if str(el.get("type") or "") != "text":
            continue
        text = str(el.get("text") or "").strip()
        bbox_pt = el.get("bbox_pt")
        if not text or not isinstance(bbox_pt, list) or len(bbox_pt) != 4:
            continue
        try:
            x0, y0, x1, y1 = [float(v) for v in bbox_pt]
        except Exception:
            continue

        px_bbox = [
            (x0 / page_w_pt) * float(width_px),
            (y0 / page_h_pt) * float(height_px),
            (x1 / page_w_pt) * float(width_px),
            (y1 / page_h_pt) * float(height_px),
        ]
        lines.append({"text": text, "bbox": px_bbox})

    return _normalize_line_items(lines)


def _load_json_lines(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return []
    if not isinstance(payload, list):
        return []
    return _normalize_line_items(payload)


@dataclass
class EngineResult:
    name: str
    status: str
    elapsed_sec: float
    info: dict[str, Any] = field(default_factory=dict)
    error: str | None = None


def _count_ppt_text_items(pptx_path: Path) -> tuple[int, int]:
    """Return (text_shape_count, non_empty_text_char_count)."""

    if not pptx_path.exists():
        return (0, 0)

    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return (0, 0)

    shape_count = 0
    char_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = str(getattr(shape, "text", "") or "").strip()
            if not text:
                continue
            shape_count += 1
            char_count += len(text)
    return (shape_count, char_count)


def _run_local_engine(
    *,
    pdf_path: Path,
    page_spans: list[tuple[int, int]],
    ocr_provider: str,
    api_key: str,
    base_url: str,
    model: str,
    ocr_strict_mode: bool,
) -> EngineResult:
    from app.job_paths import ensure_job_dir
    from app.worker import process_pdf_job

    spans = _group_contiguous_spans(
        [p for a, b in page_spans for p in range(int(a), int(b) + 1)]
    )
    if not spans:
        raise RuntimeError("empty page spans")

    chunks: list[dict[str, Any]] = []
    total_elapsed = 0.0
    total_shapes = 0
    total_chars = 0

    for span_start, span_end in spans:
        job_id = f"bench-local-{uuid.uuid4()}"
        job_dir = ensure_job_dir(job_id)
        shutil.copy2(pdf_path, job_dir / "input.pdf")

        started = time.time()
        process_pdf_job(
            job_id,
            enable_ocr=True,
            enable_layout_assist=False,
            parse_provider="local",
            page_start=int(span_start),
            page_end=int(span_end),
            ocr_provider=ocr_provider,
            ocr_ai_provider="siliconflow",
            ocr_ai_api_key=api_key,
            ocr_ai_base_url=base_url,
            ocr_ai_model=model,
            ocr_strict_mode=bool(ocr_strict_mode),
        )
        elapsed = time.time() - started
        total_elapsed += float(elapsed)

        output_pptx = job_dir / "output.pptx"
        ok = output_pptx.exists()
        if not ok:
            return EngineResult(
                name=f"local_{ocr_provider}",
                status="failed",
                elapsed_sec=total_elapsed,
                info={
                    "ocr_provider": ocr_provider,
                    "chunks": chunks,
                },
                error="output.pptx not found",
            )

        text_shapes, text_chars = _count_ppt_text_items(output_pptx)
        total_shapes += int(text_shapes)
        total_chars += int(text_chars)

        chunks.append(
            {
                "job_id": job_id,
                "job_dir": str(job_dir),
                "output_pptx": str(output_pptx),
                "page_start": int(span_start),
                "page_end": int(span_end),
            }
        )

    return EngineResult(
        name=f"local_{ocr_provider}",
        status="ok",
        elapsed_sec=total_elapsed,
        info={
            "ocr_provider": ocr_provider,
            "chunks": chunks,
            "ppt_text_shapes": int(total_shapes),
            "ppt_text_chars": int(total_chars),
        },
        error=None,
    )


def _run_v2_engine(
    *,
    pdf_path: Path,
    page_spans: list[tuple[int, int]],
    api_key: str,
    base_url: str,
    model: str,
) -> EngineResult:
    from app.job_paths import ensure_job_dir
    from app.worker import process_pdf_job

    spans = _group_contiguous_spans(
        [p for a, b in page_spans for p in range(int(a), int(b) + 1)]
    )
    if not spans:
        raise RuntimeError("empty page spans")

    chunks: list[dict[str, Any]] = []
    total_elapsed = 0.0
    total_shapes = 0
    total_chars = 0

    for span_start, span_end in spans:
        job_id = f"bench-v2-{uuid.uuid4()}"
        job_dir = ensure_job_dir(job_id)
        shutil.copy2(pdf_path, job_dir / "input.pdf")

        started = time.time()
        process_pdf_job(
            job_id,
            parse_provider="v2",
            page_start=int(span_start),
            page_end=int(span_end),
            api_key=api_key,
            base_url=base_url,
            model=model,
        )
        elapsed = time.time() - started
        total_elapsed += float(elapsed)

        output_pptx = job_dir / "output.pptx"
        ok = output_pptx.exists()
        if not ok:
            return EngineResult(
                name="v2",
                status="failed",
                elapsed_sec=total_elapsed,
                info={"chunks": chunks},
                error="output.pptx not found",
            )

        text_shapes, text_chars = _count_ppt_text_items(output_pptx)
        total_shapes += int(text_shapes)
        total_chars += int(text_chars)

        chunks.append(
            {
                "job_id": job_id,
                "job_dir": str(job_dir),
                "output_pptx": str(output_pptx),
                "page_start": int(span_start),
                "page_end": int(span_end),
            }
        )

    return EngineResult(
        name="v2",
        status="ok",
        elapsed_sec=total_elapsed,
        info={
            "chunks": chunks,
            "ppt_text_shapes": int(total_shapes),
            "ppt_text_chars": int(total_chars),
        },
        error=None,
    )


def _run_v3_engine(
    *,
    name: str,
    backend: str,
    pdf_path: Path,
    max_pages: int,
    api_key: str,
    base_url: str,
    model: str,
    out_dir: Path,
) -> EngineResult:
    # Removed: v2/v3 packages deleted
    raise RuntimeError("v3 engine support removed (ocr_pdf2ppt_v3 package deleted)")


def _safe_run(name: str, fn: Callable[[], EngineResult]) -> EngineResult:
    try:
        return fn()
    except Exception as e:
        return EngineResult(
            name=name,
            status="failed",
            elapsed_sec=0.0,
            info={},
            error=f"{type(e).__name__}: {e}\n{traceback.format_exc(limit=6)}",
        )


def _resolve_engine_page_artifacts(
    *,
    engine: EngineResult,
    page_no: int,
) -> tuple[Path | None, list[dict[str, Any]]]:
    name = engine.name
    if engine.status != "ok":
        return None, []

    def _resolve_chunk(job_kind: str) -> tuple[Path | None, int]:
        chunks = engine.info.get("chunks")
        if isinstance(chunks, list):
            for chunk in chunks:
                if not isinstance(chunk, dict):
                    continue
                try:
                    start = int(chunk.get("page_start") or 0)
                    end = int(chunk.get("page_end") or 0)
                except Exception:
                    continue
                if start <= page_no <= end and start > 0:
                    job_dir = Path(str(chunk.get("job_dir") or ""))
                    return job_dir, int(page_no - start)
        # Backward compatibility: single-job legacy payload.
        job_dir = Path(str(engine.info.get("job_dir") or ""))
        return job_dir, int(page_no - 1)

    if name.startswith("local_"):
        job_dir, idx = _resolve_chunk("local")
        if job_dir is None:
            return None, []
        clean = _find_first_existing(
            [
                job_dir / "artifacts" / "page_renders" / f"page-{idx:04d}.clean.png",
                job_dir
                / "artifacts"
                / "page_renders"
                / f"page-{idx:04d}.mineru.clean.png",
            ]
        )
        if clean is None:
            return None, []
        lines = _load_local_lines_px(
            job_dir=job_dir, page_no=page_no, clean_image=clean
        )
        return clean, lines

    if name == "v2":
        job_dir, idx = _resolve_chunk("v2")
        if job_dir is None:
            return None, []
        clean = _find_first_existing(
            [
                job_dir / "artifacts" / "page_renders" / f"page-{idx:04d}.clean.png",
                job_dir
                / "artifacts"
                / "page_renders"
                / f"page-{idx:04d}.mineru.clean.png",
            ]
        )
        if clean is None:
            return None, []
        lines = _load_local_lines_px(
            job_dir=job_dir, page_no=page_no, clean_image=clean
        )
        return clean, lines

    if name == "v2":
        job_dir, idx = _resolve_chunk("v2")
        if job_dir is None:
            return None, []
        clean = _find_first_existing(
            [
                job_dir / "artifacts" / "page_renders" / f"page-{idx:04d}.clean.png",
                job_dir
                / "artifacts"
                / "page_renders"
                / f"page-{idx:04d}.mineru.clean.png",
            ]
        )
        if clean is None:
            return None, []
        lines = _load_local_lines_px(
            job_dir=job_dir, page_no=page_no, clean_image=clean
        )
        return clean, lines

    if name in {"v3_openai_chat", "v3_paddle_doc"}:
        work_dir = Path(str(engine.info.get("work_dir") or ""))
        clean = work_dir / "clean" / f"page-{page_no:04d}.clean.png"
        lines = _load_json_lines(work_dir / "ocr" / f"page-{page_no:04d}.json")
        return (clean if clean.exists() else None), lines

    return None, []


def main() -> int:
    parser = argparse.ArgumentParser(description="Benchmark local/v2/v3 engines")
    parser.add_argument("--pdf", required=True, help="Input PDF path")
    parser.add_argument(
        "--pages",
        default="1-3",
        help="Page spec, e.g. 1-5 or 1,3,7 or random:10 or all",
    )
    parser.add_argument(
        "--seed",
        type=int,
        default=int(os.getenv("BENCH_SEED", "0") or "0"),
        help="Random seed for pages=random:N (0=use time-based seed)",
    )
    parser.add_argument(
        "--out-dir", default="../test/benchmark-engines", help="Output directory"
    )
    parser.add_argument(
        "--dpi", type=int, default=200, help="Render DPI for source image"
    )
    parser.add_argument(
        "--keep-out-dir",
        action="store_true",
        help="Keep existing out-dir content (default clears old artifacts before run)",
    )

    parser.add_argument(
        "--api-key",
        default=os.getenv("SILICONFLOW_API_KEY", ""),
        help="API key for remote OCR engines (AI OCR / PaddleOCR-VL / legacy v2 / v3). Optional for local-only runs.",
    )
    parser.add_argument(
        "--base-url",
        default=os.getenv("SILICONFLOW_BASE_URL", "https://api.siliconflow.cn/v1"),
    )
    parser.add_argument(
        "--model",
        default=os.getenv("SILICONFLOW_MODEL", "PaddlePaddle/PaddleOCR-VL-1.5"),
    )
    parser.add_argument(
        "--ocr-strict-mode",
        action="store_true",
        help="Enable strict OCR mode (disable implicit fallback providers).",
    )

    parser.add_argument(
        "--dns-host",
        default=os.getenv("BENCH_DNS_HOST", "api.siliconflow.cn"),
        help="Optional host for socket DNS override",
    )
    parser.add_argument(
        "--dns-ip",
        default=os.getenv("BENCH_DNS_IP", ""),
        help="Optional fixed IP for --dns-host",
    )

    parser.add_argument("--no-local", action="store_true", help="Skip local engine")
    parser.add_argument("--no-v2", action="store_true", help="Skip v2 engine")
    parser.add_argument(
        "--no-v3-chat", action="store_true", help="Skip v3 openai_chat engine"
    )
    parser.add_argument(
        "--no-v3-paddle-doc",
        action="store_true",
        help="Skip v3 paddle_doc_parser engine",
    )
    args = parser.parse_args()

    api_dir = _ensure_api_on_path()

    pdf_path = Path(args.pdf).resolve()
    if not pdf_path.exists():
        print(f"ERROR: pdf not found: {pdf_path}", file=sys.stderr)
        return 2

    if not args.api_key:
        print(
            "[bench] NOTE: --api-key is missing; running local OCR engines only (tesseract/paddle_local).",
            flush=True,
        )

    _clear_proxy_env()
    os.environ["REDIS_URL"] = "memory://"
    os.environ.setdefault("PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK", "True")

    restore_dns = _patch_dns(args.dns_host, (args.dns_ip or "").strip() or None)
    try:
        with pymupdf.open(str(pdf_path)) as doc:
            seed = None if int(args.seed) == 0 else int(args.seed)
            page_numbers = _parse_pages_spec(
                str(args.pages),
                max_pages=int(doc.page_count),
                seed=seed,
            )

        if not page_numbers:
            print("ERROR: no valid pages selected", file=sys.stderr)
            return 2

        out_dir = Path(args.out_dir).resolve()
        if out_dir.exists() and not args.keep_out_dir:
            shutil.rmtree(out_dir)
        out_dir = _ensure_dir(out_dir)

        page_spans = _group_contiguous_spans(page_numbers)
        page_start = min(page_numbers)
        page_end = max(page_numbers)

        engines: list[EngineResult] = []

        if not args.no_local:
            print("[bench] running local tesseract engine...", flush=True)
            engines.append(
                _safe_run(
                    "local_tesseract",
                    lambda: _run_local_engine(
                        pdf_path=pdf_path,
                        page_spans=page_spans,
                        ocr_provider="tesseract",
                        api_key=str(args.api_key or ""),
                        base_url=args.base_url,
                        model=args.model,
                        ocr_strict_mode=bool(args.ocr_strict_mode),
                    ),
                )
            )

            print("[bench] running local paddle_local engine...", flush=True)
            engines.append(
                _safe_run(
                    "local_paddle_local",
                    lambda: _run_local_engine(
                        pdf_path=pdf_path,
                        page_spans=page_spans,
                        ocr_provider="paddle_local",
                        api_key=str(args.api_key or ""),
                        base_url=args.base_url,
                        model=args.model,
                        ocr_strict_mode=bool(args.ocr_strict_mode),
                    ),
                )
            )

            if args.api_key:
                print("[bench] running local paddle (remote) engine...", flush=True)
                engines.append(
                    _safe_run(
                        "local_paddle",
                        lambda: _run_local_engine(
                            pdf_path=pdf_path,
                            page_spans=page_spans,
                            ocr_provider="paddle",
                            api_key=args.api_key,
                            base_url=args.base_url,
                            model=args.model,
                            ocr_strict_mode=bool(args.ocr_strict_mode),
                        ),
                    )
                )

                print("[bench] running local aiocr engine...", flush=True)
                engines.append(
                    _safe_run(
                        "local_aiocr",
                        lambda: _run_local_engine(
                            pdf_path=pdf_path,
                            page_spans=page_spans,
                            ocr_provider="aiocr",
                            api_key=args.api_key,
                            base_url=args.base_url,
                            model=args.model,
                            ocr_strict_mode=bool(args.ocr_strict_mode),
                        ),
                    )
                )

        if not args.no_v2:
            print("[bench] skip v2 engine (ocr_pdf2ppt_v2 package deleted)", flush=True)

        if not args.no_v3_chat:
            print(
                "[bench] skip v3 openai_chat engine (ocr_pdf2ppt_v3 package deleted)",
                flush=True,
            )

        if not args.no_v3_paddle_doc:
            print(
                "[bench] skip v3 paddle_doc_parser engine (ocr_pdf2ppt_v3 package deleted)",
                flush=True,
            )

        page_reports: list[dict[str, Any]] = []
        metrics_by_engine: dict[str, list[dict[str, float]]] = {
            engine.name: [] for engine in engines if engine.status == "ok"
        }

        for page_no in page_numbers:
            page_dir = _ensure_dir(out_dir / f"page-{page_no:03d}")
            source_path = _render_source_page(
                pdf_path,
                page_no=page_no,
                dpi=int(args.dpi),
                out_path=page_dir / "source.png",
            )
            source_img = _read_color_image(source_path)

            per_engine: dict[str, Any] = {}

            for engine in engines:
                clean_path, lines = _resolve_engine_page_artifacts(
                    engine=engine, page_no=page_no
                )
                if engine.status != "ok":
                    per_engine[engine.name] = {
                        "status": "failed",
                        "error": engine.error,
                    }
                    continue

                if clean_path is None or not clean_path.exists():
                    per_engine[engine.name] = {
                        "status": "missing_artifacts",
                        "clean_image": str(clean_path) if clean_path else None,
                        "line_count": len(lines),
                    }
                    continue

                final_path = page_dir / f"{engine.name}.final.sim.png"
                _synthesize_from_clean(
                    clean_image=clean_path, lines=lines, out_path=final_path
                )

                final_img = _align_image_like(source_img, _read_color_image(final_path))
                metrics = _image_metrics(source_img, final_img)
                metrics_by_engine.setdefault(engine.name, []).append(metrics)

                _save_side_by_side(
                    left=source_img,
                    right=final_img,
                    left_label=f"SOURCE PAGE{page_no}",
                    right_label=f"{engine.name.upper()} FINAL(SIM)",
                    out_path=page_dir / f"compare-source-vs-{engine.name}.png",
                )
                _save_heat_overlay(
                    base=source_img,
                    other=final_img,
                    out_path=page_dir / f"heat-source-vs-{engine.name}.png",
                )

                per_engine[engine.name] = {
                    "status": "ok",
                    "clean_image": str(clean_path),
                    "line_count": len(lines),
                    "final_image": str(final_path),
                    "metrics": metrics,
                }

            page_reports.append({"page_no": page_no, "engines": per_engine})

        summary_rows: list[dict[str, Any]] = []
        for engine in engines:
            text_shapes = int(engine.info.get("ppt_text_shapes") or 0)
            text_chars = int(engine.info.get("ppt_text_chars") or 0)
            row: dict[str, Any] = {
                "engine": engine.name,
                "status": engine.status,
                "elapsed_sec": engine.elapsed_sec,
                "error": engine.error,
                "info": engine.info,
                "ppt_text_shapes": text_shapes,
                "ppt_text_chars": text_chars,
            }
            series = metrics_by_engine.get(engine.name, [])
            if series:
                row["pages_scored"] = len(series)
                row["avg_mae"] = float(np.mean([m["mae"] for m in series]))
                row["avg_rmse"] = float(np.mean([m["rmse"] for m in series]))
                row["avg_high_diff_ratio"] = float(
                    np.mean([m["high_diff_ratio"] for m in series])
                )
                row["editability_bonus"] = float(
                    min(1.0, text_shapes / max(1, len(page_numbers) * 40))
                )
                text_density = float(text_chars) / float(
                    max(1, len(page_numbers) * 300)
                )
                text_quality_bonus = float(min(0.25, text_density * 0.25))
                empty_text_penalty = 0.25 if int(text_chars) <= 0 else 0.0
                row["text_density"] = text_density
                row["text_quality_bonus"] = text_quality_bonus
                row["empty_text_penalty"] = empty_text_penalty
                row["score"] = (
                    float(row["avg_mae"])
                    + 0.4 * float(row["avg_high_diff_ratio"])
                    + float(empty_text_penalty)
                    - float(text_quality_bonus)
                )
            else:
                row["pages_scored"] = 0
            summary_rows.append(row)

        rankable = [
            r
            for r in summary_rows
            if r.get("pages_scored", 0) > 0 and r.get("status") == "ok"
        ]
        rankable.sort(
            key=lambda r: (float(r["score"]), float(r["avg_mae"]), float(r["avg_rmse"]))
        )

        report = {
            "input": {
                "pdf": str(pdf_path),
                "pages": page_numbers,
                "dpi": int(args.dpi),
                "base_url": args.base_url,
                "model": args.model,
                "ocr_strict_mode": bool(args.ocr_strict_mode),
            },
            "environment": {
                "redis_url": os.getenv("REDIS_URL"),
                "paddle_disable_source_check": os.getenv(
                    "PADDLE_PDX_DISABLE_MODEL_SOURCE_CHECK"
                ),
                "dns_override": {"host": args.dns_host, "ip": args.dns_ip or None},
            },
            "summary": summary_rows,
            "ranking": rankable,
            "best_engine": rankable[0]["engine"] if rankable else None,
            "pages": page_reports,
        }

        report_path = out_dir / "report.json"
        report_path.write_text(
            json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
        )

        print(
            json.dumps(
                {
                    "report": str(report_path),
                    "best_engine": report.get("best_engine"),
                    "ranking": [
                        {
                            "engine": row["engine"],
                            "score": row.get("score"),
                            "avg_mae": row.get("avg_mae"),
                            "avg_rmse": row.get("avg_rmse"),
                            "avg_high_diff_ratio": row.get("avg_high_diff_ratio"),
                            "ppt_text_shapes": row.get("ppt_text_shapes"),
                            "ppt_text_chars": row.get("ppt_text_chars"),
                            "elapsed_sec": row.get("elapsed_sec"),
                        }
                        for row in rankable
                    ],
                },
                ensure_ascii=False,
                indent=2,
            )
        )

        return 0
    finally:
        restore_dns()


if __name__ == "__main__":
    raise SystemExit(main())
