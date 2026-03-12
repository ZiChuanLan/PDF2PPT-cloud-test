from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation


API_ROOT = Path(__file__).resolve().parents[1]
if str(API_ROOT) not in sys.path:
    sys.path.insert(0, str(API_ROOT))

from app.convert.pptx import generator


def test_visual_wrap_probe_skips_clear_single_line_label() -> None:
    should_probe = generator._should_probe_visual_wrap_for_ocr_text(
        text="Fire return intervals",
        bbox_w_pt=220.0,
        bbox_h_pt=18.0,
        baseline_ocr_h_pt=18.0,
        is_heading=False,
        wrap_hint=False,
        ocr_linebreak_assisted=False,
    )

    assert should_probe is False


def test_visual_wrap_probe_keeps_ambiguous_long_ocr_box() -> None:
    should_probe = generator._should_probe_visual_wrap_for_ocr_text(
        text=(
            "It increases the rate that nutrients become available to plants by "
            "rapidly releasing them from forest litter."
        ),
        bbox_w_pt=620.0,
        bbox_h_pt=24.0,
        baseline_ocr_h_pt=18.0,
        is_heading=False,
        wrap_hint=False,
        ocr_linebreak_assisted=False,
    )

    assert should_probe is True


def test_local_color_sampling_skips_when_upstream_color_exists() -> None:
    should_sample = generator._should_sample_local_text_colors(
        source_id="ocr",
        element_color="#1a1a1a",
    )

    assert should_sample is False


def test_page_sampling_render_skips_when_no_element_needs_probe_or_local_color() -> None:
    needs_render = generator._page_needs_ocr_sampling_render(
        page_elements=[
            {
                "source": "ocr",
                "text": "Fire return intervals",
                "bbox_pt": [20.0, 20.0, 240.0, 38.0],
                "color": "#111111",
                "ocr_linebreak_assisted": False,
            }
        ],
        page_h_pt=540.0,
        baseline_ocr_h_pt=18.0,
    )

    assert needs_render is False


def test_page_sampling_render_keeps_render_for_ambiguous_long_box() -> None:
    needs_render = generator._page_needs_ocr_sampling_render(
        page_elements=[
            {
                "source": "ocr",
                "text": (
                    "It increases the rate that nutrients become available to plants by "
                    "rapidly releasing them from forest litter."
                ),
                "bbox_pt": [20.0, 40.0, 640.0, 64.0],
                "color": "#111111",
                "ocr_linebreak_assisted": False,
            }
        ],
        page_h_pt=540.0,
        baseline_ocr_h_pt=18.0,
    )

    assert needs_render is True


def test_scanned_heading_centering_rejects_left_aligned_title() -> None:
    should_center = generator._should_center_scanned_heading(
        x0_pt=36.0,
        x1_pt=220.0,
        page_w_pt=720.0,
    )

    assert should_center is False


def test_scanned_heading_centering_accepts_centered_title() -> None:
    should_center = generator._should_center_scanned_heading(
        x0_pt=168.0,
        x1_pt=548.0,
        page_w_pt=720.0,
    )

    assert should_center is True


def test_generate_pptx_skips_final_preview_export_when_disabled(
    monkeypatch, tmp_path: Path
) -> None:
    preview_calls: list[int] = []

    monkeypatch.setattr(
        generator,
        "_export_final_preview_page_image",
        lambda **_kwargs: preview_calls.append(1),
    )

    out_path = generator.generate_pptx_from_ir(
        {
            "pages": [
                {
                    "page_index": 0,
                    "page_width_pt": 720.0,
                    "page_height_pt": 540.0,
                    "has_text_layer": True,
                    "elements": [],
                }
            ]
        },
        tmp_path / "out.pptx",
        artifacts_dir=tmp_path / "artifacts",
        export_final_preview_images=False,
    )

    assert out_path.exists()
    assert preview_calls == []


def test_generate_pptx_keeps_final_preview_export_when_enabled(
    monkeypatch, tmp_path: Path
) -> None:
    preview_calls: list[int] = []

    monkeypatch.setattr(
        generator,
        "_export_final_preview_page_image",
        lambda **_kwargs: preview_calls.append(1),
    )

    generator.generate_pptx_from_ir(
        {
            "pages": [
                {
                    "page_index": 0,
                    "page_width_pt": 720.0,
                    "page_height_pt": 540.0,
                    "has_text_layer": True,
                    "elements": [],
                }
            ]
        },
        tmp_path / "out-enabled.pptx",
        artifacts_dir=tmp_path / "artifacts-enabled",
        export_final_preview_images=True,
    )

    assert preview_calls == [1]


def test_merge_text_erase_bboxes_fast_path_keeps_transitive_line_merge() -> None:
    boxes = []
    x = 10.0
    for _ in range(260):
        boxes.append([x, 20.0, x + 8.0, 30.0])
        x += 9.0

    merged = generator._merge_text_erase_bboxes(boxes, gap_pt=2.0)

    assert merged == [[10.0, 20.0, 2349.0, 30.0]]


def test_merge_text_erase_bboxes_fast_path_keeps_separate_rows() -> None:
    boxes: list[list[float]] = []
    for y in (20.0, 44.0):
        x = 10.0
        for _ in range(130):
            boxes.append([x, y, x + 8.0, y + 10.0])
            x += 9.0

    merged = generator._merge_text_erase_bboxes(boxes, gap_pt=2.0)

    assert merged == [
        [10.0, 20.0, 1179.0, 30.0],
        [10.0, 44.0, 1179.0, 54.0],
    ]


def test_generate_pptx_reports_progress_for_scanned_pages(
    monkeypatch, tmp_path: Path
) -> None:
    def _fake_render_pdf_page_png(
        _source_pdf: Path,
        *,
        page_index: int,
        dpi: int,
        out_path: Path,
    ):
        del page_index, dpi
        img = Image.new("RGB", (800, 450), color=(255, 255, 255))
        out_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(out_path)
        return img

    monkeypatch.setattr(generator, "_render_pdf_page_png", _fake_render_pdf_page_png)
    monkeypatch.setattr(
        generator,
        "_build_scanned_image_region_infos",
        lambda **_kwargs: [],
    )

    source_pdf = tmp_path / "input.pdf"
    source_pdf.write_bytes(b"%PDF-1.4\n%stub\n")

    progress_events: list[tuple[int, int]] = []
    out_path = generator.generate_pptx_from_ir(
        {
            "source_pdf": str(source_pdf),
            "pages": [
                {
                    "page_index": 0,
                    "page_width_pt": 720.0,
                    "page_height_pt": 405.0,
                    "has_text_layer": False,
                    "elements": [
                        {
                            "type": "text",
                            "source": "ocr",
                            "text": "Fire and Ice",
                            "bbox_pt": [36.0, 48.0, 220.0, 84.0],
                            "color": "#111111",
                        }
                    ],
                }
            ],
        },
        tmp_path / "scanned-progress.pptx",
        artifacts_dir=tmp_path / "scanned-artifacts",
        export_final_preview_images=False,
        progress_callback=lambda done, total: progress_events.append((done, total)),
    )

    assert out_path.exists()
    assert progress_events == [(1, 1)]


def test_generate_pptx_scanned_heading_keeps_left_title_uncentered_and_unforced_bold(
    monkeypatch, tmp_path: Path
) -> None:
    def _fake_render_pdf_page_png(
        _source_pdf: Path,
        *,
        page_index: int,
        dpi: int,
        out_path: Path,
    ):
        del page_index, dpi
        img = Image.new("RGB", (800, 450), color=(255, 255, 255))
        out_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(out_path)
        return img

    monkeypatch.setattr(generator, "_render_pdf_page_png", _fake_render_pdf_page_png)
    monkeypatch.setattr(
        generator,
        "_build_scanned_image_region_infos",
        lambda **_kwargs: [],
    )

    source_pdf = tmp_path / "input-left-title.pdf"
    source_pdf.write_bytes(b"%PDF-1.4\n%stub\n")

    out_path = generator.generate_pptx_from_ir(
        {
            "source_pdf": str(source_pdf),
            "pages": [
                {
                    "page_index": 0,
                    "page_width_pt": 720.0,
                    "page_height_pt": 405.0,
                    "has_text_layer": False,
                    "elements": [
                        {
                            "type": "text",
                            "source": "ocr",
                            "text": "Local-first MCP",
                            "bbox_pt": [36.0, 44.0, 240.0, 82.0],
                            "color": "#111111",
                        }
                    ],
                }
            ],
        },
        tmp_path / "scanned-left-title.pptx",
        artifacts_dir=tmp_path / "scanned-left-title-artifacts",
        export_final_preview_images=False,
    )

    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    text_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and shape.text_frame is not None
        and shape.text.strip()
    ]

    assert len(text_shapes) == 1
    paragraph = text_shapes[0].text_frame.paragraphs[0]
    assert paragraph.alignment is None
    assert paragraph.runs[0].font.bold is None


def test_generate_pptx_fast_mode_skips_image_region_analysis_and_preview_export(
    monkeypatch, tmp_path: Path
) -> None:
    captured_dpis: list[int] = []
    image_region_calls = {"count": 0}
    preview_calls = {"count": 0}

    def _fake_render_pdf_page_png(
        _source_pdf: Path,
        *,
        page_index: int,
        dpi: int,
        out_path: Path,
    ):
        del page_index
        captured_dpis.append(int(dpi))
        img = Image.new("RGB", (800, 450), color=(255, 255, 255))
        out_path.parent.mkdir(parents=True, exist_ok=True)
        img.save(out_path)
        return img

    def _fake_build_scanned_image_region_infos(**_kwargs):
        image_region_calls["count"] += 1
        return []

    monkeypatch.setattr(generator, "_render_pdf_page_png", _fake_render_pdf_page_png)
    monkeypatch.setattr(
        generator,
        "_build_scanned_image_region_infos",
        _fake_build_scanned_image_region_infos,
    )
    monkeypatch.setattr(
        generator,
        "_export_final_preview_page_image",
        lambda **_kwargs: preview_calls.__setitem__("count", preview_calls["count"] + 1),
    )

    source_pdf = tmp_path / "input-fast.pdf"
    source_pdf.write_bytes(b"%PDF-1.4\n%stub\n")

    out_path = generator.generate_pptx_from_ir(
        {
            "source_pdf": str(source_pdf),
            "pages": [
                {
                    "page_index": 0,
                    "page_width_pt": 720.0,
                    "page_height_pt": 405.0,
                    "has_text_layer": False,
                    "elements": [
                        {
                            "type": "text",
                            "source": "ocr",
                            "text": "Speed first experimental mode",
                            "bbox_pt": [36.0, 48.0, 280.0, 84.0],
                            "color": "#111111",
                        }
                    ],
                }
            ],
        },
        tmp_path / "scanned-fast.pptx",
        artifacts_dir=tmp_path / "scanned-fast-artifacts",
        scanned_render_dpi=200,
        export_final_preview_images=True,
        ppt_generation_mode="fast",
    )

    assert out_path.exists()
    assert captured_dpis == [120]
    assert image_region_calls["count"] == 0
    assert preview_calls["count"] == 0
